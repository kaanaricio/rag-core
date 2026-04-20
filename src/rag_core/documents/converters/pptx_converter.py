"""PPTX converter with hybrid text extraction + OCR fallback.

Extracts text from slide shapes, tables, and speaker notes.
"""

from __future__ import annotations

import asyncio
import io
import logging
from typing import Any, Dict, List

from .base import (
    ConversionResult,
    HybridConverter,
    render_markdown_table,
    score_text_quality,
)

logger = logging.getLogger(__name__)


def _extract_shape_text(shape: Any) -> List[str]:
    """Extract text lines from a PPTX shape (text frames and tables)."""
    lines: List[str] = []

    if shape.has_text_frame:
        for paragraph in shape.text_frame.paragraphs:
            text = paragraph.text.strip()
            if text:
                lines.append(text)

    if shape.has_table:
        rows: List[List[str]] = []
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            rows.append(cells)
        if rows:
            lines.append(render_markdown_table(rows))

    return lines


def _extract_slide_figure_items(slide: Any, slide_index: int) -> List[Dict[str, Any]]:
    """Extract figure metadata from image-like slide shapes."""
    figures: List[Dict[str, Any]] = []
    try:
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore[import-not-found]
    except Exception:
        # python-pptx may be unavailable in some environments; keep extraction going.
        MSO_SHAPE_TYPE = None  # type: ignore[assignment]

    figure_number = 0
    for shape in slide.shapes:
        is_picture = False
        try:
            if MSO_SHAPE_TYPE is not None and shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                is_picture = True
        except Exception:
            # Some shapes expose shape_type inconsistently; fall back to image detection.
            pass
        if not is_picture and getattr(shape, "image", None) is not None:
            is_picture = True
        if not is_picture:
            continue

        figure_number += 1
        figure_id = "fig:slide:%d:%d" % (slide_index + 1, figure_number)
        label = "Slide %d Figure %d" % (slide_index + 1, figure_number)
        description = "Embedded image extracted from slide %d." % (slide_index + 1)
        try:
            alt_text = str(getattr(shape, "name", "") or "").strip()
            if alt_text:
                description = alt_text
        except Exception:
            # Alt text is optional; missing metadata should not block figure extraction.
            pass

        figures.append(
            {
                "figure_id": figure_id,
                "page_index": slide_index,
                "label": label,
                "description": description,
                "metadata": {
                    "source": "pptx:picture_shape",
                    "slide_number": slide_index + 1,
                },
            }
        )

    return figures


class PptxConverter(HybridConverter):
    """Converts PPTX files to markdown with slide structure and speaker notes."""

    format_name = "pptx"

    async def _try_extract(
        self,
        file_bytes: bytes,
        filename: str,
        mime_type: str,
    ) -> ConversionResult:
        """Extract text from PPTX using python-pptx."""
        from pptx import Presentation  # type: ignore[import-not-found]

        def _extract() -> ConversionResult:
            try:
                prs = Presentation(io.BytesIO(file_bytes))
            except Exception as exc:
                logger.warning("Failed to open PPTX %s: %s", filename, exc)
                return ConversionResult(
                    needs_ocr=True,
                    metadata={"parser": "local:python-pptx", "error": str(exc)},
                )

            slide_sections: List[str] = []
            figure_items: List[Dict[str, Any]] = []

            for i, slide in enumerate(prs.slides):
                parts: List[str] = ["## Slide %d" % (i + 1)]
                try:
                    title_shape = slide.shapes.title
                except Exception:
                    # Some slide masters omit a readable title placeholder; keep extraction going without it.
                    title_shape = None
                slide_title = (
                    str(getattr(title_shape, "text", "") or "").strip()
                    if title_shape is not None
                    else ""
                )
                if slide_title:
                    parts.append("### %s" % slide_title)

                for shape in slide.shapes:
                    parts.extend(_extract_shape_text(shape))

                slide_figures = _extract_slide_figure_items(slide, i)
                if slide_figures:
                    parts.append("### Figures")
                    for item in slide_figures:
                        parts.append(
                            "- **%s**: %s"
                            % (item.get("label", "Figure"), item.get("description", ""))
                        )
                    figure_items.extend(slide_figures)

                # Speaker notes
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        parts.append("\n> **Notes:** %s" % notes)

                slide_sections.append("\n\n".join(parts))

            content = "\n\n---\n\n".join(slide_sections)
            quality = score_text_quality(content)

            metadata: Dict[str, Any] = {
                "parser": "local:python-pptx",
                "slide_count": len(prs.slides),
                "needs_ocr": quality.char_count < 50,
            }
            if figure_items:
                metadata["figure_items"] = figure_items
                metadata["figure_count"] = len(figure_items)

            return ConversionResult(
                content=content,
                metadata=metadata,
                quality=quality,
            )

        return await asyncio.to_thread(_extract)
